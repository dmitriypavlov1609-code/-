#!/usr/bin/env python3
"""
Convert documents to Knowledge Base format

Converts .docx and .pptx files to markdown for KB ingestion.
Requires: python-docx, python-pptx

Install:
    pip install python-docx python-pptx
"""
import argparse
import sys
from pathlib import Path

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("Warning: python-docx not available. Install: pip install python-docx")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not available. Install: pip install python-pptx")


def convert_docx_to_markdown(docx_path: Path, output_path: Path) -> None:
    """Convert .docx to markdown."""
    if not DOCX_AVAILABLE:
        raise ImportError("python-docx not installed")

    doc = Document(docx_path)

    lines = []
    lines.append(f"# {docx_path.stem}\n")
    lines.append(f"*Конвертировано из: {docx_path.name}*\n")
    lines.append("---\n")

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Detect headings
        if para.style.name.startswith('Heading'):
            level = para.style.name.replace('Heading ', '')
            try:
                level = int(level)
                lines.append(f"{'#' * (level + 1)} {text}\n")
            except ValueError:
                lines.append(f"## {text}\n")
        else:
            lines.append(f"{text}\n")

    # Write to file
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(lines))

    print(f"✓ Converted: {docx_path.name} → {output_path.name}")


def convert_pptx_to_markdown(pptx_path: Path, output_path: Path) -> None:
    """Convert .pptx to markdown."""
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx not installed")

    prs = Presentation(pptx_path)

    lines = []
    lines.append(f"# {pptx_path.stem}\n")
    lines.append(f"*Конвертировано из: {pptx_path.name}*\n")
    lines.append("---\n")

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n## Слайд {i}\n")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Check if it's a title
                if shape == slide.shapes.title:
                    lines.append(f"### {shape.text.strip()}\n")
                else:
                    lines.append(f"{shape.text.strip()}\n")

    # Write to file
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(lines))

    print(f"✓ Converted: {pptx_path.name} → {output_path.name}")


def main():
    parser = argparse.ArgumentParser(description="Convert documents to KB markdown")
    parser.add_argument("input_file", help="Input file (.docx or .pptx)")
    parser.add_argument(
        "--output-dir",
        default="data/knowledge_base/policies",
        help="Output directory (default: data/knowledge_base/policies)",
    )

    args = parser.parse_args()

    input_path = Path(args.input_file)
    if not input_path.exists():
        print(f"Error: File not found: {input_path}")
        sys.exit(1)

    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    output_path = output_dir / f"{input_path.stem}.md"

    # Convert based on file type
    if input_path.suffix == '.docx':
        try:
            convert_docx_to_markdown(input_path, output_path)
        except ImportError:
            print("\nTo convert .docx files, install python-docx:")
            print("  pip install python-docx")
            sys.exit(1)

    elif input_path.suffix == '.pptx':
        try:
            convert_pptx_to_markdown(input_path, output_path)
        except ImportError:
            print("\nTo convert .pptx files, install python-pptx:")
            print("  pip install python-pptx")
            sys.exit(1)

    else:
        print(f"Error: Unsupported file type: {input_path.suffix}")
        print("Supported: .docx, .pptx")
        sys.exit(1)

    print(f"\n✓ Output saved to: {output_path}")
    print(f"\nNext steps:")
    print(f"  1. Review and edit {output_path}")
    print(f"  2. Run: python scripts/populate_kb.py --postgres-url ... --openai-api-key ...")


if __name__ == "__main__":
    main()
